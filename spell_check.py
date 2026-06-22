import re
import zipfile
import io
import os
from tkinter import Tk, Frame, Button, Text, Scrollbar, filedialog, messagebox, Label, simpledialog
from enchant import Dict, checker
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from xml.etree import ElementTree as ET

# 界面配色
BG_GREEN = "#CCE8CC"
ERR_SPELL_YELLOW = "#FFF280"
ERR_GRAM_BLUE = "#CCE5FF"

# 词典初始化
try:
    local_dict = Dict("en_US")
except:
    local_dict = Dict("en")
spell_checker = checker.SpellChecker(local_dict)

# ==================== 全局词汇规则 ====================
# 系动词/助动词
SING_VERBS = {"is", "was", "has", "does"}
PLUR_VERBS = {"are", "were", "have", "do"}
ALL_AUX_VERBS = SING_VERBS | PLUR_VERBS

# 情态动词
MODAL_VERBS = {"can", "could", "may", "might", "will", "would", "shall", "should", "must"}

# 合法单数/复数代词
SING_PRON = {"he", "she", "this", "that", "someone", "anyone", "everyone", "one"}
PLUR_PRON = {"we", "they", "these", "those", "both", "many", "few"}

# 黑名单（屏蔽易误判词汇）
EXCLUDE_WORDS = {
    "guide", "system", "interface", "chapter", "section", "page", "tab",
    "range", "device", "moment", "step", "sample", "problem", "air",
    "equipment", "performance", "case", "u.s.", "us",
    "in", "to", "on", "up", "into", "with", "at", "for", "by", "from",
    "still", "briefly", "downward", "no", "now", "here", "there", "also",
    "you", "each", "it", "what", "do"
}

# 冠词/限定词
DET_WORDS = {"the", "a", "an", "my", "your", "his", "her", "our", "their", "some", "any", "all"}
FULL_EXCLUDE = DET_WORDS | EXCLUDE_WORDS

# ==================== 工具函数 ====================
def is_noun_singular(word: str) -> bool:
    w = word.lower().strip()
    if w in SING_PRON:
        return True
    if w in PLUR_PRON:
        return False
    if w in FULL_EXCLUDE:
        return True
    if w.endswith(("s", "es", "ies", "ves")):
        return False
    return True

def get_nearest_noun_after_be(sent: str) -> str:
    part = re.sub(r"\s+and\s+.+", "", sent, flags=re.IGNORECASE)
    words = part.strip().split()
    for w in words:
        w_low = w.lower()
        if w_low not in FULL_EXCLUDE:
            return w
    return ""

# ==================== 语法检查核心 ====================
def check_there_be(sent: str, offset: int, full_text: str, err_list):
    pat = re.compile(r"\bthere\s+(is|are|was|were)\b", re.IGNORECASE)
    for m in pat.finditer(sent):
        verb = m.group(1).lower()
        after_be = sent[m.end():]
        nearest_noun = get_nearest_noun_after_be(after_be)
        if not nearest_noun:
            continue
        sub_sing = is_noun_singular(nearest_noun)
        verb_sing = verb in SING_VERBS
        if (sub_sing and not verb_sing) or (not sub_sing and verb_sing):
            s = offset + m.start(1)
            e = offset + m.end(1)
            err_list.append((s, e, "gram_tag", "主谓不一致"))

def check_normal_agreement(sent: str, offset: int, full_text: str, err_list):
    pat_aux = re.compile(r"\b([a-zA-Z]+)\s+(is|are|was|were|has|have|does)\b", re.IGNORECASE)
    for m in pat_aux.finditer(sent):
        sub = m.group(1).lower()
        verb = m.group(2).lower()
        if sent.lower().startswith("there "):
            continue
        if sub in FULL_EXCLUDE:
            continue
        sub_sing = is_noun_singular(sub)
        verb_sing = verb in SING_VERBS
        if (sub_sing and not verb_sing) or (not sub_sing and verb_sing):
            s = offset + m.start(2)
            e = offset + m.end(2)
            err_list.append((s, e, "gram_tag", "主谓不一致"))

    pat_pron_verb = re.compile(r"\b(he|she)\s+([a-zA-Z]+)\b", re.IGNORECASE)
    for m in pat_pron_verb.finditer(sent):
        verb = m.group(2).lower()
        if verb in ALL_AUX_VERBS or verb in MODAL_VERBS or verb in FULL_EXCLUDE:
            continue
        if not verb.endswith(("s", "es")):
            s = offset + m.start(2)
            e = offset + m.end(2)
            err_list.append((s, e, "gram_tag", "主谓不一致"))

def run_grammar(text: str, err_list):
    for m in re.finditer(r"[^.!?]+[.!?]", text):
        s_txt = m.group(0)
        s_off = m.start()
        check_there_be(s_txt, s_off, text, err_list)
        check_normal_agreement(s_txt, s_off, text, err_list)

# ==================== 主程序界面 ====================
class EnglishCheckApp:
    def __init__(self, root):
        self.root = root
        self.root.title("英文拼写语法检查器")
        self.root.geometry("1450x820")
        self.root.configure(bg=BG_GREEN)

        self.all_err_list = []
        self.cur_err_idx = -1
        self.custom_words = []

        # 顶部按钮
        top_frame = Frame(self.root, bg=BG_GREEN)
        top_frame.pack(fill="x", padx=5, pady=6)
        btn_list = [
            ("打开文件/ZIP", self.load_file),
            ("一键检查", self.full_check),
            ("上一处错误", self.prev_err),
            ("下一处错误", self.next_err),
            ("添加单词", self.add_word),
            ("词典模板", self.make_template),
            ("导入词典", self.import_dict),
            ("导出词典", self.export_dict)
        ]
        for txt, cmd in btn_list:
            Button(top_frame, text=txt, width=12, command=cmd).pack(side="left", padx=3)

        # 文本区
        edit_frame = Frame(self.root, bg=BG_GREEN)
        edit_frame.pack(fill="both", expand=True, padx=10, pady=5)
        self.text_area = Text(edit_frame, bg=BG_GREEN, font=("Consolas", 11), wrap="word")
        scroll = Scrollbar(edit_frame, command=self.text_area.yview)
        self.text_area.config(yscrollcommand=scroll.set)
        scroll.pack(side="right", fill="y")
        self.text_area.pack(fill="both", expand=True)

        # 状态栏
        self.status = Label(self.root, bg=BG_GREEN, font=("Microsoft YaHei", 10),
                            text="未加载文档 | 拼写错误:0 | 语法/格式错误:0")
        self.status.pack(anchor="w", padx=10, pady=3)

        # 标签样式
        self.text_area.tag_config("spell_tag", background=ERR_SPELL_YELLOW)
        self.text_area.tag_config("gram_tag", background=ERR_GRAM_BLUE)
        self.text_area.bind("<Button-1>", self.click_tip)

    def offset_to_pos(self, full_str, off):
        lines = full_str.splitlines(True)
        cnt = 0
        for idx, line in enumerate(lines, 1):
            if cnt + len(line) > off:
                return f"{idx}.{off - cnt}"
            cnt += len(line)
        return "end"

    def pre_clean_lines(self, text):
        lines = text.splitlines()
        res = []
        blank = 0
        for ln in lines:
            if not ln.strip():
                blank += 1
            else:
                if blank < 3:
                    res += [""] * blank
                blank = 0
                res.append(ln)
        if 0 < blank < 3:
            res += [""] * blank
        return "\n".join(res)

    def read_dita(self, f):
        try:
            tree = ET.parse(f)
            root = tree.getroot()
            buf = []
            for elem in root.iter():
                if elem.text:
                    buf.append(elem.text.strip())
            return "\n".join(buf)
        except Exception:
            return ""

    def load_file(self):
        path = filedialog.askopenfilename(
            filetypes=[
                ("支持文档", "*.docx;*.xlsx;*.pptx;*.dita;*.zip;*.txt"),
                ("文本文档", "*.txt"),
                ("Word", "*.docx"),
                ("Excel", "*.xlsx"),
                ("PPT", "*.pptx"),
                ("DITA", "*.dita"),
                ("ZIP", "*.zip")
            ]
        )
        if not path:
            return
        txt = ""
        ext = os.path.splitext(path)[1].lower()

        try:
            if ext == ".zip":
                with zipfile.ZipFile(path, "r") as zf:
                    for name in zf.namelist():
                        if name.lower().endswith(".dita") and not name.endswith("/"):
                            with zf.open(name) as f:
                                bio = io.BytesIO(f.read())
                                content = self.read_dita(bio)
                                txt += f"--- {name} ---\n{content}\n\n"
            elif ext == ".dita":
                with open(path, "rb", encoding="utf-8", errors="ignore") as f:
                    txt = self.read_dita(f)
            elif ext == ".txt":
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    txt = f.read()
            elif ext == ".docx":
                doc = Document(path)
                for para in doc.paragraphs:
                    txt += para.text + "\n"
            elif ext == ".xlsx":
                wb = load_workbook(path, read_only=True)
                ws = wb.active
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    txt += row_text + "\n"
            elif ext == ".pptx":
                ppt = Presentation(path)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt += shape.text + "\n"
        except Exception as e:
            messagebox.showerror("读取错误", f"文件读取失败：{str(e)}")
            return

        if not txt.strip():
            messagebox.showinfo("提示", "文件内无有效文本")

        txt = self.pre_clean_lines(txt)
        self.text_area.delete("1.0", "end")
        self.text_area.insert("1.0", txt)

        self.all_err_list.clear()
        self.cur_err_idx = -1
        self.status.config(text="文档已载入 | 拼写错误:0 | 语法/格式错误:0")

    def full_check(self):
        raw = self.text_area.get("1.0", "end-1c")
        self.text_area.tag_remove("spell_tag", "1.0", "end")
        self.text_area.tag_remove("gram_tag", "1.0", "end")

        # 临时列表：统一按字符偏移量收集所有错误
        temp_err = []

        # 1. 拼写错误
        spell_checker.set_text(raw)
        for err in spell_checker:
            s = err.wordpos
            e = err.wordpos + len(err.word)
            temp_err.append((s, e, "spell_tag", "拼写错误"))

        # 2. 多余空格
        for m in re.finditer(r"\s{2,}", raw):
            chunk = raw[m.start():m.end()]
            if not all(c == "\n" for c in chunk):
                temp_err.append((m.start(), m.end(), "gram_tag", "多余空格"))

        # 3. 标点后缺失空格
        for m in re.finditer(r"(?<=[a-zA-Z])([.,;:!?])(?=[a-zA-Z])", raw):
            temp_err.append((m.start(), m.end(), "gram_tag", "标点后缺失空格"))

        # 4. 标点前多余空格
        for m in re.finditer(r"\s+([.,;:!?])(?=[a-zA-Z])", raw):
            temp_err.append((m.start(), m.end(), "gram_tag", "标点前多余空格"))

        # 5. 句首大小写错误
        for m in re.finditer(r"([.!?]\s+)([a-z]\w+)", raw):
            s = m.start() + len(m.group(1))
            e = m.end()
            temp_err.append((s, e, "gram_tag", "大小写错误"))

        # 6. 主谓不一致
        run_grammar(raw, temp_err)

        # ========== 核心：按【起始字符位置】升序排序，实现从上到下依次遍历 ==========
        temp_err.sort(key=lambda x: x[0])

        # 去重（避免同一位置重复标记）
        unique_err = []
        last_s = -1
        for item in temp_err:
            s, e, tag, msg = item
            if s != last_s:
                unique_err.append(item)
                last_s = s
        self.all_err_list = unique_err

        # 统计数量
        spell_cnt = sum(1 for item in self.all_err_list if item[2] == "spell_tag")
        gram_cnt = len(self.all_err_list) - spell_cnt
        self.cur_err_idx = -1

        # 绘制所有标记
        for s, e, tag, msg in self.all_err_list:
            pos1 = self.offset_to_pos(raw, s)
            pos2 = self.offset_to_pos(raw, e)
            self.text_area.tag_add(tag, pos1, pos2)

        self.status.config(
            text=f"检查完成 | 拼写错误:{spell_cnt} | 语法/格式错误:{gram_cnt}"
        )
        if len(self.all_err_list) == 0:
            messagebox.showinfo("提示", "未发现错误")

    def next_err(self):
        if not self.all_err_list:
            messagebox.showwarning("提示", "未检测到任何错误")
            return
        self.cur_err_idx += 1
        if self.cur_err_idx >= len(self.all_err_list):
            self.cur_err_idx = 0
        s, e, tag, msg = self.all_err_list[self.cur_err_idx]
        raw = self.text_area.get("1.0", "end-1c")
        pos = self.offset_to_pos(raw, s)
        self.text_area.mark_set("insert", pos)
        self.text_area.see(pos)

    def prev_err(self):
        if not self.all_err_list:
            messagebox.showwarning("提示", "未检测到任何错误")
            return
        self.cur_err_idx -= 1
        if self.cur_err_idx < 0:
            self.cur_err_idx = len(self.all_err_list) - 1
        s, e, tag, msg = self.all_err_list[self.cur_err_idx]
        raw = self.text_area.get("1.0", "end-1c")
        pos = self.offset_to_pos(raw, s)
        self.text_area.mark_set("insert", pos)
        self.text_area.see(pos)

    def click_tip(self, event):
        pos = self.text_area.index(f"@{event.x},{event.y}")
        raw = self.text_area.get("1.0", "end-1c")
        for s, e, tag, msg in self.all_err_list:
            p1 = self.offset_to_pos(raw, s)
            p2 = self.offset_to_pos(raw, e)
            if self.text_area.compare(p1, "<=", pos) and self.text_area.compare(pos, "<", p2):
                messagebox.showinfo("错误详情", f"错误：{msg}")
                return

    def add_word(self):
        w = simpledialog.askstring("添加自定义单词", "输入单词：")
        if not w or w.strip() in self.custom_words:
            return
        w = w.strip()
        local_dict.add(w)
        self.custom_words.append(w)
        messagebox.showinfo("成功", f"已添加：{w}")

    def make_template(self):
        path = os.path.join(os.path.abspath("."), "词典模板.txt")
        with open(path, "w", encoding="utf-8") as f:
            f.write("# 每行一个单词\n")
        messagebox.showinfo("完成", f"模板已生成：{path}")

    def import_dict(self):
        path = filedialog.askopenfilename(filetypes=[("文本", "*.txt")])
        if not path:
            return
        add = 0
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line or line.startswith("#") or line in self.custom_words:
                    continue
                local_dict.add(line)
                self.custom_words.append(line)
                add += 1
        messagebox.showinfo("导入", f"新增 {add} 个单词")

    def export_dict(self):
        if not self.custom_words:
            messagebox.showinfo("提示", "暂无单词")
            return
        path = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("文本", "*.txt")])
        if not path:
            return
        with open(path, "w", encoding="utf-8") as f:
            f.write("\n".join(self.custom_words))
        messagebox.showinfo("导出完成")

if __name__ == "__main__":
    root = Tk()
    app = EnglishCheckApp(root)
    root.mainloop()