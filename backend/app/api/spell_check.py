import re
import zipfile
import io
import os
from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from docx import Document
from openpyxl import load_workbook
from xml.etree import ElementTree as ET
from app.utils.spell_checker import run_spelling_and_grammar_check, spell as runtime_spell

try:
    from pptx import Presentation
except ModuleNotFoundError:
    Presentation = None

router = APIRouter()


def _require_pptx_support():
    if Presentation is None:
        raise HTTPException(status_code=500, detail="当前环境缺少 python-pptx 依赖，暂不支持 PPTX 文件处理")

MAX_SPELL_ERRORS = 100
MAX_GRAMMAR_ERRORS = 200

SING_VERBS = {"is", "was", "has", "does"}
PLUR_VERBS = {"are", "were", "have", "do"}
ALL_AUX_VERBS = SING_VERBS | PLUR_VERBS
MODAL_VERBS = {"can", "could", "may", "might", "will", "would", "shall", "should", "must"}
SING_PRON = {"he", "she", "this", "that", "someone", "anyone", "everyone", "one"}
PLUR_PRON = {"we", "they", "these", "those", "both", "many", "few"}
EXCLUDE_WORDS = {
    "guide", "system", "interface", "chapter", "section", "page", "tab",
    "range", "device", "moment", "step", "sample", "problem", "air",
    "equipment", "performance", "case", "u.s.", "us",
    "in", "to", "on", "up", "into", "with", "at", "for", "by", "from",
    "still", "briefly", "downward", "no", "now", "here", "there", "also",
    "you", "each", "it", "what", "do"
}
FULL_EXCLUDE = ALL_AUX_VERBS | MODAL_VERBS | EXCLUDE_WORDS

LOW_LEVEL_RULES = [
    {"category": "style", "pattern": r"组份", "message": "常见错别字：建议使用“组分”"},
    {"category": "style", "pattern": r"污然", "message": "常见错别字：建议使用“污染”"},
    {"category": "style", "pattern": r"转数", "message": "术语错误：建议使用“转速”"},
    {"category": "style", "pattern": r"加入至", "message": "表达生硬：建议使用“加入”或“加至”"},
    {"category": "style", "pattern": r"用于用于", "message": "重复词：建议使用“用于”"},
    {"category": "style", "pattern": r"否则的话", "message": "口语化表达：建议使用“否则”"},
    {"category": "style", "pattern": r"非常的重要", "message": "表达错误：建议使用“非常重要”"},
    {"category": "style", "pattern": r"建议用户", "message": "表达冗余：建议使用“建议”或“请”"},
    {"category": "style", "pattern": r"仅供科研使用，不得用于临床诊断", "message": "表述冗余：建议使用“仅供科研使用”"},
    {"category": "style", "pattern": r"注意安全，注意安全事项，确保安全操作", "message": "表述重复：建议使用“注意安全操作”"},
    {"category": "style", "pattern": r"结果判读：根据说明书中的结果判读方法进行判读", "message": "表述重复：建议使用“结果判读：按说明书方法进行”"},
    {"category": "unit", "pattern": r"\b\d+\s*(?:ul|uL|Ul|UL)\b", "message": "单位格式错误：建议使用“μL”"},
    {"category": "unit", "pattern": r"\b\d+μL\b", "message": "单位格式错误：建议使用“10 μL”样式"},
    {"category": "unit", "pattern": r"\b\d+\s*mins\b", "message": "时间单位格式错误：建议使用“min”"},
    {"category": "unit", "pattern": r"\b\d+\s*rpm\b", "message": "转速格式错误：建议在数字与 rpm 之间留空格"},
    {"category": "unit", "pattern": r"\b\d+℃", "message": "温度格式错误：建议使用“°C”并统一空格"},
    {"category": "unit", "pattern": r"\b\d+\s*度\s*C\b", "message": "温度格式错误：建议使用“°C”并统一空格"},
    {"category": "unit", "pattern": r"\b\d+(?:mM|mmol/L)\b", "message": "浓度格式错误：建议使用“10 mM”样式"},
    {"category": "unit", "pattern": r"\bpH\d+(?:\.\d+)?\b", "message": "pH 格式错误：建议使用“pH 8.0”样式"},
    {"category": "style", "pattern": r"\bOD\d+\b", "message": "格式错误：建议使用“OD 260”样式"},
    {"category": "style", "pattern": r"\b\d+[xX]\b", "message": "符号格式错误：建议使用乘号“×”"},
    {"category": "style", "pattern": r"\b[A-Za-z][A-Za-z0-9]*\(", "message": "格式错误：英文缩写与括号之间建议留空格"},
    {"category": "grammar", "pattern": r"\b[a-zA-Z]\)(?=\S)", "message": "标号后缺少空格"},
    {"category": "style", "pattern": r"如下:(?=\d)", "message": "标点错误：建议使用中文冒号“：”"},
    {"category": "style", "pattern": r"强光下!", "message": "中文语境下建议使用“强光下。”"},
    {"category": "style", "pattern": r"非常重要……", "message": "标点错误：建议使用“非常重要。”"},
    {"category": "style", "pattern": r"'[^'\n\r\u0000-\u007f]*'", "message": "引号格式不统一：中文内容建议使用双引号"},
    {"category": "style", "pattern": r"图\d+", "message": "图表编号格式建议统一为“图 1”样式"},
    {"category": "style", "pattern": r"表\d+", "message": "图表编号格式建议统一为“表 1”样式"},
    {"category": "style", "pattern": r"RNA extraction kit", "message": "中英文混用：建议使用“RNA 提取试剂盒”", "applies_to": "chinese"},
    {"category": "style", "pattern": r"\bsample\b", "message": "中英文混用：建议使用“样本”", "applies_to": "chinese"},
    {"category": "style", "pattern": r"collection tube", "message": "中英文混用：建议使用“收集管”", "applies_to": "chinese"},
    {"category": "style", "pattern": r"Agarose Gel", "message": "中英文混用：建议使用“琼脂糖凝胶”", "applies_to": "chinese"},
    {"category": "style", "pattern": r"thermocycler", "message": "中英文混用：建议使用“热循环仪”", "applies_to": "chinese"},
    {"category": "style", "pattern": r"待补充|\bTBD\b", "message": "存在待补充占位：建议补全正式内容", "suggestions": ["补全正式内容"]},
]

WHITELIST_PATTERNS = {
    "product": re.compile(r"(MGISP(?:-\d+)?(?:-Smart\s+8)?|DNBSEQ(?:-[Tt]\d+[×xX]?\d+[RSrs]?)?|MGICLab(?:-FZ\d+)?|MGI)", re.IGNORECASE),
    "brand": re.compile(r"(Qubit|Eppendorf|HamiLton|Hamilton|Invitrogen|Thermo\s+Fisher\s+Scientific|BMG\s+LABTECH|AXYGEN|Greiner\s+Bio-One|Fluostar\s+Omega)", re.IGNORECASE),
    "document_id": re.compile(r"(JB-\w+-\d+|V\d+\.\d+(?:\.\d+)*|940-\d{6}-\d{2})"),
    "term": re.compile(r"(ssCir|dsDNA|PCR|RCR|DNB|DIPSEQ|OliGreen|MPC2000|ALPS\s+50V|Pos\d+~?Pos\d+|wfex|sp960)", re.IGNORECASE),
    "domain": re.compile(r"(mgi-tech\.com|global-mgitech\.com|MGI-service@mgi-tech\.com)"),
    "scientific": re.compile(r"(E\.\s*coli|in\s+situ|in\s+vitro|in\s+vivo|et\s+al\.)", re.IGNORECASE),
    "element": re.compile(r"(H|He|Li|Be|B|C|N|O|F|Ne|Na|Mg|Al|Si|P|S|Cl|Ar|K|Ca|Sc|Ti|V|Cr|Mn|Fe|Co|Ni|Cu|Zn|Ga|Ge|As|Se|Br|Kr|Rb|Sr|Y|Zr|Nb|Mo|Tc|Ru|Rh|Pd|Ag|Cd|In|Sn|Sb|Te|I|Xe|Cs|Ba|La|Ce|Pr|Nd|Pm|Sm|Eu|Gd|Tb|Dy|Ho|Er|Tm|Yb|Lu|Hf|Ta|W|Re|Os|Ir|Pt|Au|Hg|Tl|Pb|Bi|Po|At|Rn|Fr|Ra|Ac|Th|Pa|U|Np|Pu|Am|Cm|Bk|Cf|Es|Fm|Md|No|Lr)"),
}

TECH_TERMS = {
    "ng", "mg", "kg", "ug", "pg", "fg", "ml", "ul", "nl", "pl", "fl", "dl",
    "mm", "cm", "dm", "km", "nm", "um", "pm", "fm",
    "ms", "us", "ns",
    "mol", "mmol", "umol", "nmol", "pmol", "fmol",
    "sec", "min", "hr", "hrs", "mins", "secs",
    "rpm", "rcf", "od", "rbc", "wbc",
    "bp", "kb", "mb", "gb", "kda", "mda", "da",
    "atp", "nad", "fadh", "tris", "edta", "bsa", "fbs", "dmem", "rpmi", "pbs",
    "elisa", "chip", "chipseq", "facs", "western", "blot", "wb", "ip",
    "illumina", "nanopore", "pacbio", "bgi", "qiagen", "invitrogen",
    "miseq", "hiseq", "nextseq", "novaseq", "minion", "promethion",
    "rnase", "dnase", "hek", "cho", "gfp", "yfp", "cfp", "rfp", "egfp",
    "ph", "coli",
    "qpcr", "rtpcr", "ngs", "wgs", "wes", "wts",
    "dna", "rna", "rrna", "mrna", "scrna", "cdna", "pcr", "ssdna", "dsdna", "gdna",
    "dnb", "dnbseq", "mgi", "dnbe", "dnbs", "omics", "pipetting", "biosafety",
    "extranet", "thermo", "oligo", "rxn", "qubit", "milli-q", "dnbseqtm", "qubittm",
    "mgisp", "mgiclab", "qubit", "eppendorf", "hamilton", "hamiilton",
    "bmglabtech", "axygen", "greiner", "fluostar", "omega",
    "sscir", "rcr", "dipseq", "oligreen", "mpc2000", "alps",
    "wfex", "sp960",
    "situ", "vitro", "vivo",
    "coli", "tech", "mgitech", "hs", "xte", "te",

    # 技术术语（从CSV提取）
    "abdomen",
    "aberration",
    "abort",
    "aborted",
    "accessories",
    "accuracy",
    "adapter",
    "adenine",
    "adenovirus",
    "admin",
    "aerosols",
    "agilent",
    "airdetectionsensitivity",
    "airgapdetection",
    "alarm",
    "allele",
    "alpha",
    "amphibians",
    "analyze",
    "antibody",
    "aperture",
    "api",
    "apoptosis",
    "app",
    "appendix",
    "approver",
    "approx",
    "apptitle",
    "archaea",
    "aspirate",
    "aspirateoffsetdirection",
    "aspirator",
    "assembling",
    "assess",
    "async",
    "auto",
    "autosome",
    "autostep",
    "avg",
    "bacteria",
    "baffle",
    "bai",
    "barcode",
    "barcodefilepath",
    "barcodelength",
    "base",
    "basecall",
    "basecallip",
    "baseline",
    "baygene",
    "beaker",
    "beta",
    "binuclease",
    "bioinformatics",
    "bioligo",
    "biosafety",
    "bioverse",
    "blang",
    "bonan",
    "braille",
    "bubbles",
    "buckle",
    "buret",
    "buyei",
    "buzzer",
    "calchannelcounts",
    "calibration",
    "carotid",
    "carp",
    "cartridgeinplace",
    "caster",
    "cat",
    "ccd",
    "centrifuge",
    "chastity",
    "chloroplasts",
    "chromosome",
    "chuck",
    "cine",
    "clamping",
    "class",
    "cleanfcscriptpath",
    "cleanliness",
    "cloning",
    "clotdetection",
    "clotdetectionsensitivity",
    "cnt",
    "coagulation",
    "codon",
    "collapse",
    "colony",
    "commissioning",
    "compatibility",
    "compatible",
    "compensation",
    "compensations",
    "complaint",
    "concentrator",
    "condenser",
    "conductive",
    "cone",
    "confidential",
    "config",
    "conflicting",
    "conformity",
    "congenital",
    "constraint",
    "consumable",
    "consumables",
    "contig",
    "contraindications",
    "controlboard",
    "conversion",
    "coturnserveraddr",
    "counterclockwise",
    "creep",
    "crosstalk",
    "crystal",
    "crystallization",
    "css",
    "culture",
    "cyclization",
    "cycloneseq",
    "cytoactivity",
    "cytolysis",
    "cytosine",
    "dai",
    "daur",
    "deang",
    "debris",
    "degeneracy",
    "dehumidifier",
    "deletion",
    "delta",
    "demo",
    "demulsification",
    "denaturation",
    "denature",
    "denoise",
    "derung",
    "desiccator",
    "designer",
    "detailed",
    "dia",
    "diaphragm",
    "dicot",
    "dideoxynucleotides",
    "digestion",
    "dipeptide",
    "direction",
    "disconnect",
    "disconnecting",
    "disinfectant",
    "dislodge",
    "dispense",
    "dispenseheight",
    "dispensemanner",
    "distributor",
    "dntp",
    "dog",
    "dominant",
    "dong",
    "dongle",
    "dongxiang",
    "downsampling",
    "drawer",
    "driver",
    "droplets",
    "dsdna",
    "duck",
    "duplication",
    "durability",
    "eclinkrfid",
    "electrophoresis",
    "elute",
    "emptyairoffsetofz",
    "emptyairrate",
    "endexam",
    "engineer",
    "english",
    "enhancer",
    "enrich",
    "epe",
    "epi",
    "equilibrate",
    "etc",
    "ethanol",
    "euchromatin",
    "eui",
    "evaporator",
    "ewenki",
    "exam",
    "examconclusion",
    "execute",
    "exit",
    "exiting",
    "exome",
    "exon",
    "expand",
    "expandpos",
    "extend",
    "extendtype",
    "eyepiece",
    "factory",
    "failed",
    "farther",
    "female",
    "fence",
    "figure",
    "filtered",
    "final",
    "firstspeedofz",
    "fixation",
    "fixture",
    "flatness",
    "flowcellinplace",
    "fluorometer",
    "fluorophore",
    "foci",
    "folk",
    "forceps",
    "formamide",
    "fragmentation",
    "freezer",
    "ftat",
    "fungi",
    "funnel",
    "furnace",
    "gaoshan",
    "gasket",
    "gdna",
    "gelao",
    "gender",
    "gene",
    "genewell",
    "genotype",
    "genotyping",
    "genus",
    "gin",
    "goat",
    "goose",
    "grasp",
    "grind",
    "gripper",
    "guanine",
    "guest",
    "gui",
    "gynecology",
    "hamster",
    "han",
    "hani",
    "hanshine",
    "haploid",
    "haploinsufficiency",
    "harm",
    "hazard",
    "heart",
    "hemocytometer",
    "hemolysis",
    "hemophilia",
    "heterochromatin",
    "heterozygous",
    "hikvision",
    "hollow",
    "homepage",
    "homogenize",
    "homologs",
    "homology",
    "homozygous",
    "hotplate",
    "hour",
    "housing",
    "html",
    "hui",
    "humidity",
    "hybridization",
    "hydrophilic",
    "hydrophobic",
    "idcard",
    "identity",
    "idx",
    "ignore",
    "illustration",
    "imageformat",
    "imager",
    "imageregion",
    "immunotherapy",
    "importer",
    "importing",
    "impurities",
    "inactivate",
    "incubate",
    "info",
    "ingroup",
    "inherited",
    "inhibitor",
    "initialization",
    "initialize",
    "initializing",
    "integrity",
    "interaction",
    "introduction",
    "intron",
    "invertebrate",
    "io",
    "irradiance",
    "isextend",
    "isexternalexc",
    "ismultipoint",
    "isochores",
    "isopropanol",
    "isopropyl",
    "issync",
    "isverifydata",
    "ivd",
    "jingpo",
    "jino",
    "joystick",
    "js",
    "kazak",
    "keyword",
    "kidney",
    "kilobase",
    "kingdom",
    "kirgiz",
    "labeling",
    "labware",
    "lahu",
    "landscape",
    "lane",
    "laser",
    "leukemia",
    "lever",
    "lhoba",
    "libraryid",
    "lid",
    "lifecycle",
    "lint",
    "liquidclass",
    "lisu",
    "lithography",
    "live",
    "liver",
    "locus",
    "loglevel",
    "loop",
    "loopindex",
    "lot",
    "lung",
    "lymphocyte",
    "lysozyme",
    "macroparticle",
    "magnification",
    "male",
    "malformation",
    "mammals",
    "man",
    "manage",
    "manifests",
    "manifold",
    "manufacture",
    "maonan",
    "mapping",
    "marking",
    "mask",
    "max",
    "melanoma",
    "melway",
    "membership",
    "metagenomics",
    "metaphase",
    "metatranscriptomics",
    "methylation",
    "miao",
    "microarray",
    "microchannel",
    "micrometer",
    "microorganisms",
    "micropipette",
    "microplate",
    "microsatellite",
    "microswitch",
    "microtome",
    "microvibration",
    "mindetectoffsetofz",
    "mindvision",
    "minisatellite",
    "misalignment",
    "misfocusing",
    "mission",
    "mitochondrion",
    "mixaspirateheight",
    "mixaspirateoffset",
    "mixdispenseheight",
    "mixdispenseoffset",
    "mixdispenseoffsetdirection",
    "mixemptyheight",
    "mixemptyoffset",
    "mixemptyoffsetdirection",
    "mixture",
    "modifying",
    "monba",
    "mongol",
    "monocot",
    "monosomy",
    "mortar",
    "motif",
    "msg",
    "mulao",
    "multimeter",
    "multipointinfo",
    "mutation",
    "mute",
    "naxi",
    "neurofibromatosis",
    "neuroscience",
    "never",
    "newtemplate",
    "noise",
    "normal",
    "normalization",
    "nsclc",
    "nuclease",
    "octagonal",
    "ocular",
    "offline",
    "offsetspeedofxy",
    "oligonucleotide",
    "oncogene",
    "open",
    "openpore",
    "operational",
    "operon",
    "optical",
    "optocoupler",
    "oroqen",
    "orthology",
    "os",
    "oscilloscope",
    "others",
    "outgroup",
    "outpatient",
    "oven",
    "overview",
    "pack",
    "paraconfig",
    "parafilm",
    "parallel",
    "paralogy",
    "parasite",
    "pass",
    "pathogen",
    "paused",
    "pausing",
    "pdf",
    "pedigree",
    "pellet",
    "peltier",
    "pending",
    "permeabilization",
    "personnel",
    "phantom",
    "pharynx",
    "phylogenetic",
    "phylum",
    "pickupoffsetz",
    "pierce",
    "pipette",
    "plant",
    "plasma",
    "play",
    "plunger",
    "polydactyly",
    "pooling",
    "portrait",
    "positioning",
    "positivesign",
    "postairrate",
    "poultry",
    "powder",
    "preairrate",
    "precipitates",
    "preface",
    "prepare",
    "preparing",
    "pressed",
    "preventdrop",
    "preview",
    "previous",
    "primer",
    "printer",
    "probe",
    "profilometer",
    "promoter",
    "pronucleus",
    "proposedvalue",
    "protein",
    "proteome",
    "proteomics",
    "prototyping",
    "pulse",
    "pumi",
    "pump",
    "purpose",
    "pushreport",
    "putdownoffsetz",
    "qiang",
    "quantification",
    "rabbit",
    "rack",
    "radiance",
    "radiator",
    "rat",
    "read",
    "reading",
    "ready",
    "real",
    "recessive",
    "recipe",
    "recipes",
    "red",
    "reflectance",
    "reflector",
    "refresh",
    "refrigerator",
    "registration",
    "relay",
    "reload",
    "remarks",
    "remoteaddress",
    "rename",
    "renaturation",
    "renishaw",
    "replace",
    "reserved",
    "reservoir",
    "reset",
    "reside",
    "residual",
    "restart",
    "resume",
    "resuspend",
    "retrovirus",
    "retry",
    "return",
    "rinse",
    "rod",
    "roomid",
    "roycom",
    "russia",
    "salar",
    "sata",
    "scaffold",
    "scan",
    "scanning",
    "scope",
    "script",
    "scrofa",
    "sdk",
    "sealing",
    "secondspeedofz",
    "secret",
    "segments",
    "selected",
    "sending",
    "sensitivity",
    "sent",
    "sequencer",
    "sequencing",
    "serum",
    "servertitle",
    "settempduration",
    "settimeout",
    "several",
    "severity",
    "shaker",
    "she",
    "sheep",
    "shell",
    "shutdown",
    "signalserveraddr",
    "silencer",
    "simulate",
    "simulated",
    "sipper",
    "site",
    "skirted",
    "slider",
    "softdelete",
    "solid",
    "sonographer",
    "speaker",
    "species",
    "specificity",
    "specimen",
    "spleen",
    "ssdna",
    "ssn",
    "stack",
    "startchar",
    "starteventlog",
    "starting",
    "startmicro",
    "startspeaker",
    "statement",
    "station",
    "status",
    "stepping",
    "stirrer",
    "stopcock",
    "stopeventlog",
    "stopmicro",
    "stopping",
    "stopspeaker",
    "streaks",
    "stringconst",
    "subcellular",
    "subheading",
    "subject",
    "submit",
    "substitution",
    "substrate",
    "subtyping",
    "successfully",
    "sui",
    "sum",
    "supernatant",
    "supplier",
    "swab",
    "switch",
    "syringe",
    "tajik",
    "tatar",
    "temp",
    "temperatureduration",
    "terms",
    "thaw",
    "threadlocker",
    "thymine",
    "thyroid",
    "tile",
    "timeout",
    "tips",
    "title",
    "tmp",
    "toggle",
    "tonsils",
    "toolbar",
    "tooling",
    "torque",
    "trained",
    "transcription",
    "transgene",
    "transgenic",
    "translocation",
    "translucent",
    "transmittance",
    "tray",
    "trisomy",
    "trusee",
    "trypsin",
    "ts",
    "tujia",
    "tumor",
    "typing",
    "uhrr",
    "ui",
    "unclassified",
    "uniformity",
    "unknown",
    "unspecified",
    "unused",
    "uracil",
    "username",
    "ux",
    "uygur",
    "uzbek",
    "vacuole",
    "variants",
    "vector",
    "verification",
    "verifydata",
    "verifytype",
    "vertebrate",
    "vial",
    "viscosity",
    "visibility",
    "vision",
    "voltmeter",
    "vortex",
    "vs",
    "waiting",
    "wash",
    "washing",
    "web",
    "well",
    "worklist",
    "workmode",
    "workpiece",
    "xibe",
    "yao",
    "yes",
    "yugur",
    "zang",
    "zebracall",
    "zebrafish",
    "zhuang",
}

RE_WORDS = re.compile(r"[a-zA-Z]+")
RE_MULTI_SPACE = re.compile(r"\s{2,}")
RE_MISSING_SPACE_AFTER_PUNCT = re.compile(r"(?<=[a-zA-Z])([.,;:!?])(?=[a-zA-Z])")
RE_EXTRA_SPACE_BEFORE_PUNCT = re.compile(r"\s+([.,;:!?])(?=[a-zA-Z])")
RE_CASE_ERROR = re.compile(r"([.!?]\s+)([a-z]\w+)")
RE_SENTENCE = re.compile(r"[^.!?]+[.!?]")
RE_THERE_BE = re.compile(r"\bthere\s+(is|are|was|were)\b", re.IGNORECASE)
RE_AGREEMENT = re.compile(r"\b([a-zA-Z]+)\s+(is|are|was|were|has|have|does)\b", re.IGNORECASE)
RE_PRON_VERB = re.compile(r"\b(he|she)\s+([a-zA-Z]+)\b", re.IGNORECASE)
# 中文标点混入英文
RE_CHINESE_PUNCT = re.compile(r"([a-zA-Z])、([a-zA-Z])")
RE_CHINESE_DOT = re.compile(r"(?<=[a-zA-Z])。(?=[a-zA-Z])")

# 单位检查规则
RE_MICROLITER = re.compile(r"(?<=\d)ul\b|(?<=\d)uL\b")
RE_MILLILITER = re.compile(r"(?<=\d)ml\b")
RE_TIME_UNIT = re.compile(r"(?<=\d)mins\b|(?<=\d)hs\b|(?<=\d)sec\b")
RE_UNIT_CASE = re.compile(r"\b(Kg|CM|KM|ML|MG|UG|PG|FG|DL|DM|NM|PM|FM|US)\b")
RE_NUMBER_UNIT_SPACE = re.compile(r"\d+[μμmMkK][Ll]")

# 乘号检测（检测x/X代替×）
RE_MULTIPLY_SYMBOL = re.compile(r"(\d+)[xX]([A-Za-z]+)")

# 单词拆分检测
RE_WORD_SPLIT = re.compile(r"\b(desk)\s+(top)\b|\b(back)\s+(up)\b|\b(set)\s+(up)\b|\b(front)\s+(end)\b|\b(right)\s+(click)\b|\b(left)\s+(click)\b|\b(enter)\s+(key)\b|\b(delete)\s+(key)\b")

# 官网地址检查
RE_OLD_DOMAIN = re.compile(r"en\.mgi-tech\.com")

# 联系方式检查
RE_WRONG_EMAIL = re.compile(r"[a-zA-Z0-9._%+-]+@(?!mgi-tech\.com)[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}")

# 风格建议模式
RE_STYLE_PATTERNS = [
    (re.compile(r"\bPlease\b", re.IGNORECASE), "Please出现在正文", "建议使用更中性的表达", "suggestion"),
    (re.compile(r"\bat\s+your\s+own\s+risk\b", re.IGNORECASE), "at your own risk", "考虑使用更正式的表述", "suggestion"),
    (re.compile(r"\bat\s+their\s+own\s+risk\b", re.IGNORECASE), "at their own risk", "考虑使用更正式的表述", "suggestion"),
]

# 语法错误模式
RE_GRAMMAR_PATTERNS = [
    (re.compile(r"\bfor\s+run\s+\w+\b", re.IGNORECASE), "for run one plate", "建议改为 for running one plate 或 for one run of", "suggestion"),
]


def is_word_in_whitelist(word, text, start, end):
    for pattern in WHITELIST_PATTERNS.values():
        if pattern.search(text, start, end + 1):
            return True
    return False


def _parse_issue_position(position):
    try:
        start_text, end_text = str(position or '').split('-', 1)
        start = int(start_text)
        end = int(end_text)
        if start < 0 or end < start:
            raise ValueError
        return start, end
    except (ValueError, TypeError, AttributeError):
        return None, None


def _extract_suggestions(issue):
    suggestion = str(issue.get('suggestion') or '').strip()
    if not suggestion:
        return []
    parts = re.split(r'[、,;\n]+', suggestion)
    cleaned = []
    seen = set()
    for part in parts:
        value = part.strip()
        if not value or value in seen:
            continue
        seen.add(value)
        cleaned.append(value)
    return cleaned[:5]


def _extract_message_suggestions(message):
    candidates = re.findall(r'“([^”]+)”', str(message or ''))
    cleaned = []
    seen = set()
    for candidate in candidates:
        value = candidate.strip()
        if not value or value in seen:
            continue
        seen.add(value)
        cleaned.append(value)
    return cleaned[:3]


def _build_low_level_suggestions(rule, matched_text):
    explicit = rule.get('suggestions') or []
    if explicit:
        return explicit[:3]

    if rule['pattern'] == r'\b\d+\s*(?:ul|uL|Ul|UL)\b':
        number = re.search(r'\d+', matched_text)
        return [f"{number.group(0)} μL"] if number else ['μL']
    if rule['pattern'] == r'\b\d+\s*mins\b':
        number = re.search(r'\d+', matched_text)
        return [f"{number.group(0)} min"] if number else ['min']
    if rule['pattern'] == r'\b\d+\s*rpm\b':
        return [re.sub(r'(?i)rpm$', ' rpm', matched_text)]
    if rule['pattern'] == r'\b\d+μL\b':
        return [re.sub(r'(\d+)μL', r'\1 μL', matched_text)]
    if rule['pattern'] == r'\b\d+℃':
        return [matched_text.replace('℃', ' °C')]
    if rule['pattern'] == r'\b\d+\s*度\s*C\b':
        number = re.search(r'\d+', matched_text)
        return [f"{number.group(0)} °C"] if number else []
    if rule['pattern'] == r'\b\d+(?:mM|mmol/L)\b':
        return [re.sub(r'^(\d+)', r'\1 ', matched_text)]
    if rule['pattern'] == r'\bpH\d+(?:\.\d+)?\b':
        return [matched_text.replace('pH', 'pH ', 1)]
    if rule['pattern'] == r'\bOD\d+\b':
        return [matched_text.replace('OD', 'OD ', 1)]
    if rule['pattern'] == r'\b\d+[xX]\b':
        return [re.sub(r'[xX]$', '×', matched_text)]
    if rule['pattern'] == r'\b[A-Za-z][A-Za-z0-9]*\(': 
        return [matched_text[:-1] + ' (']
    if rule['pattern'] == r'\b[a-zA-Z]\)(?=\S)':
        return [matched_text + ' ']
    if rule['pattern'] == r'如下:(?=\d)':
        return [matched_text.replace(':', '：', 1)]
    if rule['pattern'] == r"'[^'\n\r\u0000-\u007f]*'":
        return [matched_text.replace("'", '"')]
    if rule['pattern'] == r'图\d+':
        return [matched_text[0] + ' ' + matched_text[1:]]
    if rule['pattern'] == r'表\d+':
        return [matched_text[0] + ' ' + matched_text[1:]]

    message_suggestions = _extract_message_suggestions(rule.get('message'))
    if message_suggestions:
        return message_suggestions
    return []


def _append_match_issues(issues, text, pattern, *, category, message, suggestion, source):
    for match in re.finditer(pattern, text, re.IGNORECASE):
        issues.append({
            'severity': 'general',
            'category': category,
            'source': source,
            'original_text': text[match.start():match.end()],
            'context': text[max(0, match.start() - 60):min(len(text), match.end() + 60)],
            'description': message,
            'suggestion': suggestion,
            'position': f"{match.start()}-{match.end()}",
        })


def _normalize_date_suggestion(matched_text, document_language):
    parts = re.findall(r'\d+', matched_text)
    if len(parts) != 3:
        return '07/27/2026' if document_language == 'english' else '2026-07-27'

    year, month, day = parts
    if document_language == 'english':
        return f"{month.zfill(2)}/{day.zfill(2)}/{year.zfill(4)}"
    return f"{year.zfill(4)}-{month.zfill(2)}-{day.zfill(2)}"


def _collect_consistency_issues(text, document_language):
    issues = []

    if 'Cat.No' in text and 'Cat. No.' in text:
        _append_match_issues(
            issues,
            text,
            r'Cat\.\s?No\.?|Catalog Number',
            category='style',
            message='术语格式混用：货号名称请统一使用“Cat. No.”或“Catalog Number”中的一种',
            suggestion='Cat. No.、Catalog Number',
            source='consistency_rule',
        )

    if 'PBMC' in text and '外周血单个核细胞' in text:
        _append_match_issues(
            issues,
            text,
            r'PBMC|外周血单个核细胞',
            category='style',
            message='术语混用：请统一使用“PBMC”或“外周血单个核细胞”中的一种',
            suggestion='PBMC、外周血单个核细胞',
            source='consistency_rule',
        )

    if 'RNA 完整性' in text and 'RNA 质量' in text:
        _append_match_issues(
            issues,
            text,
            r'RNA\s*完整性|RNA\s*质量',
            category='style',
            message='术语混用：请统一使用“RNA 完整性”或“RNA 质量”中的一种',
            suggestion='RNA 完整性、RNA 质量',
            source='consistency_rule',
        )

    numbering_patterns = [r'\b\d+\.', r'\b\d+\)', r'\(\d+\)', r'①|②|③|④|⑤|⑥|⑦|⑧|⑨|⑩']
    for line_match in re.finditer(r'.*', text):
        line = line_match.group(0)
        if not line.strip():
            continue
        matched_numbering_styles = sum(bool(re.search(pattern, line)) for pattern in numbering_patterns)
        if matched_numbering_styles < 2:
            continue
        line_offset = line_match.start()
        for pattern in numbering_patterns:
            for match in re.finditer(pattern, line):
                issues.append({
                    'severity': 'general',
                    'category': 'style',
                    'source': 'consistency_rule',
                    'original_text': match.group(0),
                    'context': line,
                    'description': '编号格式混用：请统一全文编号层级样式',
                    'suggestion': '1.',
                    'position': f"{line_offset + match.start()}-{line_offset + match.end()}",
                })

    date_patterns = [
        r'\b\d{4}-\d{2}-\d{2}\b',
        r'\b\d{4}/\d{2}/\d{2}\b',
        r'\b\d{4}\.\d{2}\.\d{2}\b',
    ]
    matched_date_styles = sum(bool(re.search(pattern, text)) for pattern in date_patterns)
    if matched_date_styles >= 2:
        for pattern in date_patterns:
            for match in re.finditer(pattern, text):
                matched_text = match.group(0)
                issues.append({
                    'severity': 'general',
                    'category': 'style',
                    'source': 'consistency_rule',
                    'original_text': matched_text,
                    'context': text[max(0, match.start() - 60):min(len(text), match.end() + 60)],
                    'description': '日期格式混用：请统一全文日期格式',
                    'suggestion': _normalize_date_suggestion(matched_text, document_language),
                    'position': f"{match.start()}-{match.end()}",
                })

    if 'min、分钟、mins' in text or ('分钟' in text and re.search(r'\bmins?\b', text, re.IGNORECASE)):
        _append_match_issues(
            issues,
            text,
            r'\bmins?\b|分钟',
            category='style',
            message='时间单位混用：请统一使用“min”或“分钟”中的一种',
            suggestion='min、分钟',
            source='consistency_rule',
        )

    if re.search(r'\bDNBSEQ\b', text) and re.search(r'\bDNBseq\b', text):
        _append_match_issues(
            issues,
            text,
            r'\bDNBSEQ\b|\bDNBseq\b',
            category='style',
            message='产品名称大小写混用：请统一使用“DNBSEQ”或“DNBseq”中的一种',
            suggestion='DNBSEQ、DNBseq',
            source='consistency_rule',
        )

    return issues


def _map_issue_type(issue):
    source = str(issue.get('source') or '').lower()
    category = str(issue.get('category') or '')
    category_lower = category.lower()
    if source == 'spellcheck' or '拼写' in category:
        return 'spell'
    if category_lower == 'unit' or '单位' in category:
        return 'unit'
    if category_lower == 'style' or '风格' in category or '格式' in category:
        return 'style'
    if category_lower == 'grammar':
        return 'grammar'
    return 'grammar'


def _build_response(text, issues):
    errors = []
    seen = set()
    for issue in issues:
        start, end = _parse_issue_position(issue.get('position'))
        if start is None or end is None or end > len(text):
            continue
        issue_type = _map_issue_type(issue)
        dedupe_key = (start, end, issue_type)
        if dedupe_key in seen:
            continue
        seen.add(dedupe_key)
        word = issue.get('original_text') or text[start:end]
        errors.append({
            'start': start,
            'end': end,
            'type': issue_type,
            'severity': issue.get('severity') or 'general',
            'message': issue.get('description') or issue.get('category') or '发现可疑问题',
            'word': word,
            'context': issue.get('context') or '',
            'suggestions': _extract_suggestions(issue),
        })

    errors.sort(key=lambda item: (item['start'], item['end'], item['type']))
    spell_count = sum(1 for item in errors if item['type'] == 'spell')
    grammar_count = len(errors) - spell_count
    return {
        'errors': errors,
        'spell_count': spell_count,
        'grammar_count': grammar_count,
        'total_count': len(errors),
        'text': text,
    }


def _detect_document_language(text):
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
    english_words = len(re.findall(r'\b[a-zA-Z]{2,}\b', text))
    if english_words >= 5 and english_words >= max(1, chinese_chars * 2):
        return 'english'
    return 'chinese'


def _collect_low_level_rule_issues(text, document_language):
    issues = []
    for rule in LOW_LEVEL_RULES:
        applies_to = rule.get('applies_to')
        if applies_to and applies_to != document_language:
            continue
        pattern = re.compile(rule['pattern'], re.IGNORECASE)
        for match in pattern.finditer(text):
            matched_text = text[match.start():match.end()]
            issues.append({
                'severity': 'general',
                'category': rule['category'],
                'source': 'low_level_rule',
                'original_text': matched_text,
                'context': text[max(0, match.start() - 50):min(len(text), match.end() + 50)],
                'description': rule['message'],
                'suggestion': '、'.join(_build_low_level_suggestions(rule, matched_text)),
                'position': f"{match.start()}-{match.end()}",
            })
    return issues


def process_text(text):
    """共享处理函数：统一走完整规则链并适配前端结果结构"""
    normalized_text = pre_clean_lines(text)
    document_language = _detect_document_language(normalized_text)
    issues = run_spelling_and_grammar_check(normalized_text)
    issues.extend(_collect_low_level_rule_issues(normalized_text, document_language))
    issues.extend(_collect_consistency_issues(normalized_text, document_language))
    return _build_response(normalized_text, issues)


def is_noun_singular(word: str) -> bool:
    w = word.lower().strip()
    if w in SING_PRON:
        return True
    if w in PLUR_PRON:
        return False
    if w in FULL_EXCLUDE:
        return True
    if w.endswith(("s", "es", "ies", "ves")):
        return False
    return True


def get_nearest_noun_after_be(sent: str) -> str:
    part = re.sub(r"\s+and\s+.+", "", sent, flags=re.IGNORECASE)
    words = part.strip().split()
    for w in words:
        w_low = w.lower()
        if w_low not in FULL_EXCLUDE:
            return w
    return ""


def check_there_be(sent: str, offset: int, full_text: str, err_list):
    for m in RE_THERE_BE.finditer(sent):
        verb = m.group(1).lower()
        after_be = sent[m.end():]
        nearest_noun = get_nearest_noun_after_be(after_be)
        if not nearest_noun:
            continue
        sub_sing = is_noun_singular(nearest_noun)
        verb_sing = verb in SING_VERBS
        if (sub_sing and not verb_sing) or (not sub_sing and verb_sing):
            s = offset + m.start(1)
            e = offset + m.end(1)
            err_list.append({"start": s, "end": e, "type": "grammar", "severity": "general", "message": "主谓不一致"})


def check_normal_agreement(sent: str, offset: int, full_text: str, err_list):
    for m in RE_AGREEMENT.finditer(sent):
        sub = m.group(1).lower()
        verb = m.group(2).lower()
        if sent.lower().startswith("there "):
            continue
        if sub in FULL_EXCLUDE:
            continue
        sub_sing = is_noun_singular(sub)
        verb_sing = verb in SING_VERBS
        if (sub_sing and not verb_sing) or (not sub_sing and verb_sing):
            s = offset + m.start(2)
            e = offset + m.end(2)
            err_list.append({"start": s, "end": e, "type": "grammar", "severity": "general", "message": "主谓不一致"})

    for m in RE_PRON_VERB.finditer(sent):
        verb = m.group(2).lower()
        if verb in ALL_AUX_VERBS or verb in MODAL_VERBS or verb in FULL_EXCLUDE:
            continue
        if not verb.endswith(("s", "es")):
            s = offset + m.start(2)
            e = offset + m.end(2)
            err_list.append({"start": s, "end": e, "type": "grammar", "severity": "general", "message": "主谓不一致"})


def run_grammar(text: str, err_list):
    for m in RE_SENTENCE.finditer(text):
        s_txt = m.group(0)
        s_off = m.start()
        check_there_be(s_txt, s_off, text, err_list)
        check_normal_agreement(s_txt, s_off, text, err_list)


def read_dita(f):
    try:
        tree = ET.parse(f)
        root = tree.getroot()
        buf = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                buf.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                buf.append(elem.tail.strip())
        return "\n".join(buf)
    except Exception:
        return ""


def read_idml(f):
    try:
        buf = []
        with zipfile.ZipFile(f, "r") as zf:
            for name in zf.namelist():
                name_lower = name.lower()
                if not (name_lower.startswith("stories/") and name_lower.endswith(".xml")):
                    continue
                try:
                    with zf.open(name) as xml_file:
                        tree = ET.parse(xml_file)
                        root = tree.getroot()
                        current_parts = []

                        def flush_line():
                            if not current_parts:
                                return
                            line = "".join(current_parts)
                            line = re.sub(r"\s+", " ", line).strip()
                            current_parts.clear()
                            if line:
                                buf.append(line)

                        for elem in root.iter():
                            tag = elem.tag.rsplit('}', 1)[-1] if isinstance(elem.tag, str) else ''
                            if tag == 'ParagraphStyleRange':
                                flush_line()
                                continue
                            if tag == 'Br':
                                flush_line()
                                continue
                            if tag != 'Content':
                                continue
                            if elem.text:
                                current_parts.append(elem.text)
                        flush_line()
                except Exception:
                    continue
        return "\n".join(buf)
    except Exception:
        return ""


def pre_clean_lines(text):
    lines = text.splitlines()
    res = []
    blank = 0
    for ln in lines:
        if not ln.strip():
            blank += 1
        else:
            if blank < 3:
                res += [""] * blank
            blank = 0
            res.append(ln)
    if 0 < blank < 3:
        res += [""] * blank
    return "\n".join(res)


def _iter_docx_blocks(doc):
    from docx.document import Document as DocxDocumentClass
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    parent = doc.element.body if isinstance(doc, DocxDocumentClass) else doc
    for child in parent.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def _extract_text_from_docx_document(doc):
    parts = []
    for block in _iter_docx_blocks(doc):
        rows = getattr(block, "rows", None)
        if rows is not None:
            for row in rows:
                cells = [re.sub(r"\s+", " ", cell.text or "").strip() for cell in row.cells]
                cells = [cell for cell in cells if cell]
                if cells:
                    parts.append("\t".join(cells))
            continue

        text = re.sub(r"\s+", " ", getattr(block, "text", "") or "").strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def extract_text_from_file(file: UploadFile):
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    content = file.file.read()

    if ext == ".zip":
        txt = ""
        with zipfile.ZipFile(io.BytesIO(content), "r") as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                name_lower = name.lower()
                try:
                    with zf.open(name) as f:
                        file_content = f.read()
                        if name_lower.endswith(".dita") or name_lower.endswith(".xml"):
                            bio = io.BytesIO(file_content)
                            dita_content = read_dita(bio)
                            if dita_content.strip():
                                txt += f"--- {name} ---\n{dita_content}\n\n"
                        elif name_lower.endswith(".md"):
                            md_content = file_content.decode('utf-8', errors='ignore')
                            if md_content.strip():
                                txt += f"--- {name} ---\n{md_content}\n\n"
                        elif name_lower.endswith(".idml"):
                            bio = io.BytesIO(file_content)
                            idml_content = read_idml(bio)
                            if idml_content.strip():
                                txt += f"--- {name} ---\n{idml_content}\n\n"
                        elif name_lower.endswith(".txt"):
                            txt += f"--- {name} ---\n{file_content.decode('utf-8', errors='ignore')}\n\n"
                        elif name_lower.endswith(".docx"):
                            bio = io.BytesIO(file_content)
                            doc = Document(bio)
                            doc_text = _extract_text_from_docx_document(doc)
                            if doc_text.strip():
                                txt += f"--- {name} ---\n{doc_text}\n\n"
                        elif name_lower.endswith(".xlsx"):
                            bio = io.BytesIO(file_content)
                            wb = load_workbook(bio, read_only=True)
                            ws = wb.active
                            xlsx_text = ""
                            for row in ws.iter_rows(values_only=True):
                                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                                if row_text.strip():
                                    xlsx_text += row_text + "\n"
                            if xlsx_text.strip():
                                txt += f"--- {name} ---\n{xlsx_text}\n\n"
                        elif name_lower.endswith(".pptx"):
                            _require_pptx_support()
                            bio = io.BytesIO(file_content)
                            prs = Presentation(bio)
                            pptx_text = ""
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        pptx_text += shape.text + "\n"
                            if pptx_text.strip():
                                txt += f"--- {name} ---\n{pptx_text}\n\n"
                except Exception as e:
                    continue
        return txt
    elif ext == ".dita" or ext == ".xml":
        bio = io.BytesIO(content)
        return read_dita(bio)
    elif ext == ".md":
        return content.decode("utf-8", errors="ignore")
    elif ext == ".idml":
        bio = io.BytesIO(content)
        return read_idml(bio)
    elif ext == ".txt":
        return content.decode("utf-8", errors="ignore")
    elif ext == ".docx":
        bio = io.BytesIO(content)
        doc = Document(bio)
        return _extract_text_from_docx_document(doc)
    elif ext == ".xlsx":
        bio = io.BytesIO(content)
        wb = load_workbook(bio, read_only=True)
        ws = wb.active
        txt = ""
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            txt += row_text + "\n"
        return txt
    elif ext == ".pptx":
        _require_pptx_support()
        bio = io.BytesIO(content)
        ppt = Presentation(bio)
        txt = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt += shape.text + "\n"
        return txt
    else:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")


class SpellCheckRequest(BaseModel):
    text: str


@router.post("/check", summary="检查文本拼写语法")
async def check_spell(request: SpellCheckRequest):
    if not request.text.strip():
        return {"errors": [], "spell_count": 0, "grammar_count": 0, "total_count": 0, "text": ""}
    return process_text(request.text)


@router.post("/upload", summary="上传文件并检查拼写语法")
async def upload_and_check(file: UploadFile = File(...)):
    try:
        text = extract_text_from_file(file)
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件读取失败: {str(e)}")

    if not text.strip():
        return {"errors": [], "spell_count": 0, "grammar_count": 0, "total_count": 0, "text": "", "filename": file.filename}

    result = process_text(text)
    result["filename"] = file.filename
    return result


@router.post("/add-word", summary="添加自定义单词到词典")
async def add_custom_word(word: str):
    word = word.strip()
    if not word:
        raise HTTPException(status_code=400, detail="单词不能为空")
    runtime_spell.word_frequency.load_words([word])
    return {"message": f"已添加单词: {word}"}


@router.post("/import-dict", summary="导入词典文件")
async def import_dict(file: UploadFile = File(...)):
    content = await file.read()
    lines = content.decode("utf-8").splitlines()
    words_to_add = []
    for line in lines:
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        words_to_add.append(line)
    runtime_spell.word_frequency.load_words(words_to_add)
    return {"message": f"成功导入 {len(words_to_add)} 个单词"}


@router.get("/export-dict", summary="导出自定义词典")
async def export_dict():
    words = []
    for w in runtime_spell.word_frequency:
        words.append(w)
    return {"words": words}
